import streamlit as st
import os
import json
import pandas as pd
from pptx import Presentation
from openai import OpenAI
import shutil

# -----------------------------
# Setup (Cloud)
# -----------------------------
if "OPENAI_API_KEY" not in st.secrets:
    st.error("❌ No API key found. Please add it in Streamlit Cloud → Settings → Secrets.")
    st.stop()

client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

DATA_DIR = "data"
TEMPLATE = "templates/BIP_MCG_Case Study_Insert Case Study Name.pptx"

# -----------------------------
# Helpers
# -----------------------------
def extract_text_from_ppt(ppt_file):
    """Extract all text from a PPTX file."""
    prs = Presentation(ppt_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return text_runs


def save_extracted(case_name, texts, original_file):
    """Save raw.json, parsed.csv, and original PPT in a case folder."""
    folder = os.path.join(DATA_DIR, case_name)
    os.makedirs(folder, exist_ok=True)

    # Save JSON
    with open(os.path.join(folder, "raw.json"), "w", encoding="utf-8") as f:
        json.dump({"text": texts}, f, indent=2)

    # Save CSV
    df = pd.DataFrame({"content": texts})
    df.to_csv(os.path.join(folder, "parsed.csv"), index=False)

    # Save original PPT with original name
    orig_name = f"Original - {original_file.name}"
    with open(os.path.join(folder, orig_name), "wb") as f:
        f.write(original_file.getbuffer())

    return folder


def analyze_case(texts):
    """Call OpenAI to parse case study into structured JSON."""
    joined_text = " ".join(texts)[:8000]  # truncate if huge
    prompt = f"""
    You are a consultant. Read this case study text and return structured JSON.
    Rules:
    - Never include real client names: replace with "the client".
    - Always refer to delivering company as "BIP".
    - Case Study Name:
      • If the PPT text has a clear subject (e.g., Structured Notes), use it directly.
      • If not, infer a professional name based on the content.
      • Keep it short: max 3–4 words, no dashes/colons.
    - Category: must be a consulting subcategory (e.g., Regulatory Compliance, Data Migration, Risk & Controls).
      • Never just "Financial Services" or another generic industry.
    - Function: which office(s) are impacted (e.g., COO Office, PMO Office).
    - Challenge, Solution, Results:
      • 3–4 sentences each.
      • Concise enough to fit within 3–4 lines in a PPT textbox.
      • Do not exceed 4 sentences.
    - Business Categories: 5 unique items, ≤2 words each, no repeats.
    - Hashtags: 3 unique buzzwords about products/solutions/core areas,
      • 2–3 words each,
      • no repeats,
      • return WITHOUT the # symbol (because # will be a separate text box in PPT).

    Case Study Text:
    {joined_text}

    Return JSON with keys:
    case_study_name, category, function, challenge, solution, results, business_categories, hashtags
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"}
    )

    return response.choices[0].message.content


def replace_in_shape(shape, placeholder, value):
    """Replace placeholder text in runs, preserving formatting."""
    if not shape.has_text_frame or not placeholder:
        return False
    replaced = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if placeholder in r.text:
                r.text = r.text.replace(placeholder, value)
                replaced = True
    return replaced


def create_case_ppt(analysis_json, folder):
    """Generate a new PPT from the template with parsed fields (fonts preserved)."""
    prs = Presentation(TEMPLATE)
    data = json.loads(analysis_json)

    # -----------------------------
    # Post-processing
    # -----------------------------
    name = data.get("case_study_name", "").strip()
    name = name.replace("(", "").replace(")", "")
    data["case_study_name"] = " ".join(name.split()[:6]) or "Case Study"

    for field in ["challenge", "solution", "results"]:
        text_val = data.get(field, "")
        sentences = text_val.split(". ")
        if len(sentences) > 4:
            text_val = ". ".join(sentences[:4])
        data[field] = text_val.strip()

    hashtags = list(dict.fromkeys(data.get("hashtags") or []))[:3]
    cleaned_tags = []
    for tag in hashtags:
        tag = tag.replace("#", "").strip()
        tag = " ".join(tag.split()[:3])
        if tag:
            cleaned_tags.append(tag)
    while len(cleaned_tags) < 3:
        cleaned_tags.append("")
    hashtags = cleaned_tags

    bcs = data.get("business_categories") or []
    seen = set()
    clean_bcs = []
    for x in bcs:
        item = " ".join(x.split()[:2]).strip()
        if item and item.lower() not in seen:
            seen.add(item.lower())
            clean_bcs.append(item)
    while len(clean_bcs) < 5:
        clean_bcs.append("")
    clean_bcs = clean_bcs[:5]

    # Replace in template
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, "Insert name of case study", data["case_study_name"])
            replace_in_shape(shape, "Insert Category", data.get("category", ""))
            replace_in_shape(shape, "Insert Function", data.get("function", ""))
            replace_in_shape(shape, "Insert Challenge Here", data.get("challenge", ""))
            replace_in_shape(shape, "Insert Solution Here", data.get("solution", ""))
            replace_in_shape(shape, "Insert Results Here", data.get("results", ""))
            replace_in_shape(shape, "TagOne", hashtags[0])
            replace_in_shape(shape, "TagTwo", hashtags[1])
            replace_in_shape(shape, "TagThree", hashtags[2])
            replace_in_shape(shape, "Business Category 1", clean_bcs[0])
            replace_in_shape(shape, "Business Category 2", clean_bcs[1])
            replace_in_shape(shape, "Business Category 3", clean_bcs[2])
            replace_in_shape(shape, "Business Category 4", clean_bcs[3])
            replace_in_shape(shape, "Business Category 5", clean_bcs[4])

    out_file = os.path.join(folder, f"BIP_MCG_Case Study_{data['case_study_name']}.pptx")
    prs.save(out_file)
    return out_file


def create_zip(folder, case_name):
    """Package the case study folder into a ZIP."""
    zip_path = shutil.make_archive(folder, "zip", root_dir=folder)
    return zip_path

# -----------------------------
# Streamlit App (Cloud)
# -----------------------------
st.title("📊 BIP Case Study Processor (Cloud)")

uploaded_files = st.file_uploader(
    "Upload up to 20 PPTX files",
    accept_multiple_files=True,
    type=["pptx"]
)

if uploaded_files:
    for f in uploaded_files[:20]:  # cap at 20
        case_name = os.path.splitext(f.name)[0]
        st.write(f"Processing {case_name}...")

        # 1. Extract text
        texts = extract_text_from_ppt(f)

        # 2. Save JSON + CSV + original PPT
        folder = save_extracted(case_name, texts, f)

        # 3. Analyze with OpenAI
        analysis_json = analyze_case(texts)

        # 4. Generate new PPT
        out_ppt = create_case_ppt(analysis_json, folder)

        # 5. Package into ZIP for download
        zip_path = create_zip(folder, case_name)
        with open(zip_path, "rb") as zf:
            st.download_button(
                label=f"⬇️ Download {case_name} ZIP",
                data=zf.read(),
                file_name=f"{case_name}_case_study.zip",
                mime="application/zip"
            )

        st.success(f"✅ Case study {case_name} packaged!")

